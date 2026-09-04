#!/usr/bin/env python3
"""Render the video's slides from the F5 corporate template.

The cards were drawn by tools/card.swift, which is fine as furniture but is not
F5's type, colour or logo. This builds them as real slides on the corporate
template instead, so the video opens and closes the way an F5 deck does.

The template is NOT vendored here. It is an internal asset — the repo it comes
from says so explicitly and is deliberately unpublished — and bnkscope-field is
public. It is read from wherever --template points, the built .pptx and the
rendered PNGs land in build/, and build/ is git-ignored.

    ./make-slides.py --template ../../why-ctl-tools/templates/F5-...pptx \
                     --out build/slides
"""
import argparse
import os
import subprocess
import sys
import tempfile

from pptx import Presentation
from pptx.util import Pt

# name -> (layout, title, subtitle). The names are what build-video.sh asks for.
SLIDES = [
    ("title", "Title Slide: Black Solid", "bnkscope Field",
     "Kubernetes and F5 TMM telemetry, from a Mac"),
    ("close", "Title Slide: Black Solid", "github.com/mwiget/bnkscope-field",
     "Notarized macOS download · builds for iPad from source"),
]


def layout(prs, name):
    found = {}
    for master in prs.slide_masters:
        for l in master.slide_layouts:
            found.setdefault(l.name, l)
    if name in found:
        return found[name]
    for key, value in found.items():
        if key.lower().startswith(name.lower()):
            return value
    sys.exit(f"no layout matching {name!r}; have: {sorted(found)[:8]}…")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True)
    ap.add_argument("--out", default="build/slides")
    ap.add_argument("--dpi", type=int, default=144)   # 13.333in * 144 = 1920px
    args = ap.parse_args()

    prs = Presentation(args.template)
    # The template ships with example slides; the deck is built from layouts.
    for i in range(len(prs.slides) - 1, -1, -1):
        rid = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rid)
        del prs.slides._sldIdLst[i]

    for _, layout_name, title, subtitle in SLIDES:
        slide = prs.slides.add_slide(layout(prs, layout_name))
        for ph in list(slide.placeholders):
            idx = ph.placeholder_format.idx
            if idx == 0:
                ph.text_frame.text = title
            elif idx in (11, 17):
                ph.text_frame.text = subtitle
            else:
                ph._element.getparent().remove(ph._element)

    os.makedirs(args.out, exist_ok=True)
    with tempfile.TemporaryDirectory() as work:
        pptx = os.path.join(work, "slides.pptx")
        prs.save(pptx)
        subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                        "--outdir", work, pptx], check=True,
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        pdf = os.path.join(work, "slides.pdf")
        subprocess.run(["pdftoppm", "-r", str(args.dpi), "-png", pdf,
                        os.path.join(work, "s")], check=True)
        pages = sorted(f for f in os.listdir(work) if f.startswith("s-") and f.endswith(".png"))
        if len(pages) != len(SLIDES):
            sys.exit(f"rendered {len(pages)} pages for {len(SLIDES)} slides")
        for (name, *_), page in zip(SLIDES, pages):
            dest = os.path.join(args.out, name + ".png")
            os.replace(os.path.join(work, page), dest)
            print(f"  {name:8} {dest}")


if __name__ == "__main__":
    main()
